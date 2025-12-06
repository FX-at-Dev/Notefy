# backend/python-worker/main.py
from fastapi import FastAPI, File, UploadFile, Form
from pptx import Presentation
from pypdf import PdfReader
from fastapi.responses import JSONResponse
import tempfile
import base64
import io

app = FastAPI()

@app.post("/parse-pdf")
async def parse_pdf(file: UploadFile = File(...)):
    contents = await file.read()
    pdf_file = io.BytesIO(contents)
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n\n"
    
    return JSONResponse({"text": text})

@app.post("/parse-pptx")
async def parse_pptx(file: UploadFile = File(...), mode: str = Form('slides')):
    # Save temp file
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    contents = await file.read()
    tmp.write(contents)
    tmp.flush()
    prs = Presentation(tmp.name)
    slides_out = []
    for slide in prs.slides:
        text_runs = []
        images = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_runs.append(shape.text)
            if 'picture' in shape.__class__.__name__.lower():
                # extract image - python-pptx extraction is a bit involved
                try:
                    img = shape.image
                    img_bytes = img.blob
                    b64 = base64.b64encode(img_bytes).decode('utf-8')
                    images.append(f"data:{img.content_type};base64,{b64}")
                except Exception:
                    pass
        slides_out.append({
            "text": "\n".join(text_runs).strip(),
            "images": images
        })
    return JSONResponse({"slides": slides_out})
